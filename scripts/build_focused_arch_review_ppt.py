#!/usr/bin/env python3
"""生成聚焦推荐链路的 FILA 架构评审 PPT。

风格参考: docs/架构变更评审.pptx
内容来源: docs/FILA+穿搭推荐系统架构.pdf 与现有代码
范围限定: ETL、索引、意图解析、搭配召回、LLM 排序和理由生成
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "docs" / "架构变更评审.pptx"
OUTPUT = ROOT / "docs" / "FILA穿搭推荐系统架构评审_推荐链路聚焦.pptx"

RED = RGBColor(0xC0, 0x00, 0x00)
DARK_RED = RGBColor(0x88, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x59, 0x59, 0x59)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
LIGHT_RED = RGBColor(0xFF, 0xE1, 0xE1)
PALE_RED = RGBColor(0xFB, 0xEA, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD0, 0xD0, 0xD0)
GREEN = RGBColor(0x2E, 0x8B, 0x57)

FONT = "微软雅黑"


def remove_all_slides(prs: Presentation) -> None:
    """删除模板页，仅保留母版和布局。"""
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        rel_id = slide_id.get(qn("r:id"))
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def add_slide(prs: Presentation, layout_name: str):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_layouts[0])


def set_font(run, size=12, bold=False, color=DARK) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    font_props = run._r.get_or_add_rPr()
    ea = font_props.find(qn("a:ea"))
    if ea is None:
        ea = font_props.makeelement(qn("a:ea"), {})
        font_props.append(ea)
    ea.set("typeface", FONT)


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 12,
    bold: bool = False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = anchor
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return box


def add_title(slide, text: str):
    return add_text(
        slide,
        text,
        0.45,
        0.22,
        12.4,
        0.68,
        size=28,
        bold=True,
        color=RED,
    )


def add_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill=LIGHT_GRAY,
    line=LINE,
    line_width=0.75,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_label(slide, text: str, left: float, top: float, width: float):
    shape = add_box(
        slide,
        left,
        top,
        width,
        0.34,
        fill=RED,
        line=RED,
        line_width=0,
    )
    frame = shape.text_frame
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    set_font(run, size=11, bold=True, color=WHITE)
    return shape


def add_box_text(
    slide,
    title: str,
    body: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill=WHITE,
    title_color=RED,
    body_size=10,
):
    add_box(slide, left, top, width, height, fill=fill)
    add_text(
        slide,
        title,
        left + 0.08,
        top + 0.08,
        width - 0.16,
        0.28,
        size=12,
        bold=True,
        color=title_color,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        body,
        left + 0.12,
        top + 0.44,
        width - 0.24,
        height - 0.5,
        size=body_size,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )


def add_bullets(
    slide,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 11,
):
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = Pt(4)
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = DARK
        for run in paragraph.runs:
            set_font(run, size=size, color=DARK)
    return box


def add_table(slide, headers: list[str], rows: list[list[str]], box):
    left, top, width, height = box
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = table_shape.table
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RED
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                set_font(run, size=10, bold=True, color=WHITE)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = PALE_RED if row_idx % 2 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    set_font(run, size=9, color=DARK)
    return table


def cover(prs: Presentation) -> None:
    slide = add_slide(prs, "3_浅色风格封面")
    add_text(
        slide,
        "技术架构评审",
        0.8,
        2.45,
        11.7,
        1.0,
        size=54,
        bold=True,
        color=RED,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "FILA 穿搭推荐系统",
        0.8,
        3.55,
        11.7,
        0.6,
        size=28,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "聚焦：数据 ETL · 索引构建 · 意图解析 · 搭配召回 · LLM 排序与理由",
        0.8,
        4.28,
        11.7,
        0.45,
        size=16,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "2026.07",
        0.8,
        5.08,
        11.7,
        0.35,
        size=14,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )


def scope_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "1_浅色风格封面")
    add_title(slide, "评审范围与主结论")
    cols = [
        ("离线数据层", "Hive 日更 CSV -> skus.jsonl / 固定搭配 JSON"),
        ("在线索引层", "ES 单品/搭配索引 + Milvus 图文/文本向量"),
        ("推荐编排层", "RecommendService 串联意图、召回、排序、理由"),
        ("模型增强层", "Intent LLM、Ranking LLM、Reason LLM"),
    ]
    for idx, (title, body) in enumerate(cols):
        add_box_text(
            slide,
            title,
            body,
            0.55 + idx * 3.08,
            1.28,
            2.82,
            1.35,
            fill=LIGHT_RED,
            body_size=10,
        )
    add_box(slide, 0.55, 2.95, 6.0, 3.25, fill=LIGHT_GRAY)
    add_label(slide, "本次包含", 0.7, 3.1, 1.25)
    add_bullets(
        slide,
        [
            "ETL：商品目录、图片索引图、固定搭配 JSON、增量状态。",
            "索引：ES skus/outfits，Milvus sku_vectors / sku_text_vectors。",
            "线上：意图融合、锚点识别、多路搭配召回、融合去重。",
            "模型：粗排截断、LLM 精排、搭配级/单品级推荐理由。",
        ],
        0.75,
        3.65,
        5.5,
        2.2,
        size=12,
    )
    add_box(slide, 6.85, 2.95, 6.0, 3.25, fill=LIGHT_GRAY)
    add_label(slide, "边界说明", 7.0, 3.1, 1.25)
    add_bullets(
        slide,
        [
            "页面只围绕推荐引擎主链路展开，不扩展外围展示能力。",
            "模型能力只覆盖意图、排序和理由三个在线决策点。",
            "工程章节只保留与数据、索引、召回质量直接相关的内容。",
            "单品互补能力仅作为召回背景，不单独展开评审。",
        ],
        7.05,
        3.65,
        5.45,
        2.2,
        size=12,
    )
    add_text(
        slide,
        "主结论：当前系统是一条以 ES + Milvus 为在线索引、以 LLM "
        "增强理解和排序的单服务推荐流水线，核心风险集中在数据一致性、"
        "召回路径可控性和 LLM 成本/延迟。",
        0.55,
        6.55,
        12.2,
        0.5,
        size=12,
        color=GRAY,
    )


def etl_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "1_浅色风格封面")
    add_title(slide, "数据 ETL 与日更链路")
    steps = [
        ("1 Hive 日更", "daily_download_product_tables.py\nproduct_* / guide_* CSV"),
        ("2 商品目录", "build_catalog.py\nskus.jsonl / spu_to_skus"),
        ("3 图片字段", "select_images.py\ndisplay / index images"),
        ("4 固定搭配", "build_fila_guide_outfits_fast.py\nfila_outfits.json"),
        ("5 索引写入", "build_fila_es_index.py\nMilvus index scripts"),
    ]
    for idx, (title, body) in enumerate(steps):
        left = 0.48 + idx * 2.55
        add_box_text(slide, title, body, left, 1.35, 2.3, 1.35, fill=LIGHT_RED)
        if idx < len(steps) - 1:
            add_text(
                slide,
                "->",
                left + 2.26,
                1.78,
                0.25,
                0.3,
                size=16,
                bold=True,
                color=RED,
                align=PP_ALIGN.CENTER,
            )
    add_box(slide, 0.55, 3.05, 5.9, 2.55, fill=LIGHT_GRAY)
    add_label(slide, "产物", 0.72, 3.2, 0.9)
    add_bullets(
        slide,
        [
            "data/tables：Hive 日更商品、图片、搭配源表。",
            "data/processed/skus.jsonl：在线 SKU 主数据和检索文案。",
            "data/processed/spu_to_skus.json：款号到 SKU 聚合。",
            "data/preview/fila_outfits.json：统一固定搭配数据。",
        ],
        0.75,
        3.65,
        5.45,
        1.6,
        size=11,
    )
    add_box(slide, 6.85, 3.05, 6.0, 2.55, fill=LIGHT_GRAY)
    add_label(slide, "增量机制", 7.02, 3.2, 1.2)
    add_bullets(
        slide,
        [
            "daily_incremental_update.py 统一编排下载、ETL、ES、Milvus。",
            "默认增量：product_master.updated_at + 索引状态文件。",
            "data/logs/fila_index_sync_state.json 管理文档哈希和向量签名。",
            "支持 --full 全量重建，支持 --prune-orphans 清理孤立文档。",
        ],
        7.05,
        3.65,
        5.45,
        1.6,
        size=11,
    )
    add_text(
        slide,
        "评审关注：ETL 是线上召回质量的源头，需保证 SKU 字段、角色、"
        "中类、色系、季节、索引图片与固定搭配关系的同步一致。",
        0.55,
        6.15,
        12.2,
        0.45,
        size=12,
        color=GRAY,
    )


def index_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "1_浅色风格封面")
    add_title(slide, "索引构建与在线存储")
    headers = ["类型", "对象", "数据源", "在线用途"]
    rows = [
        [
            "ES",
            "umalog-q-maiamgs-index-fila-skus",
            "data/processed/skus.jsonl",
            "Query2ES 属性/文本召回、SKU 详情 mget",
        ],
        [
            "ES",
            "umalog-q-maiamgs-index-fila-outfits",
            "fila_outfits.json + 运营搭配源",
            "固定搭配反查、按 SKU 批量查搭配",
        ],
        [
            "Milvus",
            "fila_sku_vectors",
            "SKU index_images 多图向量",
            "上传图/锚点图近邻，group_by sku_id",
        ],
        [
            "Milvus",
            "fila_sku_text_vectors",
            "SKU title + 色系 + 季节文本",
            "文本语义召回，按 role/gender/season expr 过滤",
        ],
        [
            "Milvus",
            "fila_sku_complementary_vectors",
            "互补模型向量",
            "互补拼套通路，可按开关启用",
        ],
    ]
    table = add_table(slide, headers, rows, (0.45, 1.2, 12.45, 3.0))
    for idx, width in enumerate([1.2, 3.1, 3.1, 5.05]):
        table.columns[idx].width = Inches(width)
    add_box(slide, 0.55, 4.55, 5.9, 1.65, fill=LIGHT_GRAY)
    add_label(slide, "ES 构建", 0.72, 4.72, 1.1)
    add_bullets(
        slide,
        [
            "build_fila_es_index.py 创建 skus/outfits mapping。",
            "同 _id 覆盖写入，增量模式基于 doc_hash 跳过未变文档。",
            "DataFacade 运行时 ES 优先，本地 JSON fallback。",
        ],
        0.75,
        5.1,
        5.45,
        0.85,
        size=10,
    )
    add_box(slide, 6.85, 4.55, 6.0, 1.65, fill=LIGHT_GRAY)
    add_label(slide, "Milvus 构建", 7.02, 4.72, 1.25)
    add_bullets(
        slide,
        [
            "图向量字段 product_vector，文本向量字段 text_vector。",
            "向量维度 1024，距离 COSINE，云端 HNSW，本地自动切换。",
            "增量签名覆盖图片 URL、文本、维度与 embedding 模型。",
        ],
        7.05,
        5.1,
        5.45,
        0.85,
        size=10,
    )
    add_text(
        slide,
        "评审关注：ES 承接结构化检索与固定搭配反查，Milvus 承接图文"
        "语义近邻；二者依赖同一份 skus.jsonl，字段漂移会直接影响召回。",
        0.55,
        6.45,
        12.2,
        0.45,
        size=12,
        color=GRAY,
    )


def online_flow_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "1_浅色风格封面")
    add_title(slide, "线上推荐主链路")
    stages = [
        ("用户输入", "文本 / 图片 / SKU"),
        ("意图解析", "parse_user_intent\nIntentResult"),
        ("锚点识别", "显式 SKU\n图向量近邻"),
        ("搭配召回", "固定搭配\n语义拼套\nQuery2ES"),
        ("排序理由", "粗排截断\nLLM 精排\nReason"),
        ("卡片输出", "outfit_card\noutfit_results"),
    ]
    for idx, (title, body) in enumerate(stages):
        left = 0.45 + idx * 2.08
        add_box_text(slide, title, body, left, 1.35, 1.86, 1.45, fill=LIGHT_RED)
        if idx < len(stages) - 1:
            add_text(
                slide,
                "->",
                left + 1.83,
                1.85,
                0.25,
                0.3,
                size=15,
                bold=True,
                color=RED,
                align=PP_ALIGN.CENTER,
            )
    add_box(slide, 0.55, 3.25, 12.3, 2.55, fill=LIGHT_GRAY)
    add_label(slide, "核心编排", 0.72, 3.42, 1.25)
    add_bullets(
        slide,
        [
            "RecommendService 初始化 LocalDataStore、DataFacade、SkuRetriever，"
            "在线数据读取统一经 DataFacade。",
            "chat_stream 通过 SSE 暴露阶段事件，但本评审只关注服务内数据流，"
            "不展开页面交互和辅助诊断事件。",
            "多路召回返回候选套装后先 merge_and_dedupe_outfits，"
            "再 coarse_rank_outfits 截断，最终 rank_deduped_outfits 精排。",
            "输出由 outfit_card 组装，保留 outfit_id、recall_source、items、"
            "reason、rank_score 等业务字段。",
        ],
        0.75,
        3.9,
        11.8,
        1.45,
        size=11,
    )
    add_text(
        slide,
        "评审关注：线上主链路是同步编排的模型增强检索链路，端到端延迟"
        "主要由 Milvus/ES 网络 IO、LLM 排序与理由生成决定。",
        0.55,
        6.32,
        12.2,
        0.45,
        size=12,
        color=GRAY,
    )


def intent_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "1_浅色风格封面")
    add_title(slide, "线上意图解析与锚点融合")
    add_box(slide, 0.55, 1.2, 3.85, 4.85, fill=LIGHT_GRAY)
    add_label(slide, "输入源", 0.72, 1.38, 1.05)
    add_bullets(
        slide,
        [
            "文本 Query：自然语言、货号、场景、颜色、季节、预算。",
            "selected_sku_id：用户已锁定单品，直接注入权威属性。",
            "图片：base64 embedding -> Milvus 近邻，辅助识别锚点角色。",
        ],
        0.75,
        1.85,
        3.35,
        1.25,
        size=11,
    )
    add_label(slide, "解析组件", 0.72, 3.45, 1.25)
    add_bullets(
        slide,
        [
            "query_understanding.parse_user_intent",
            "intent_engine.extract_intent",
            "llm_client.extract_intent_json",
            "role_resolver.resolve_roles",
        ],
        0.75,
        3.92,
        3.35,
        1.1,
        size=10,
    )
    add_box(slide, 4.75, 1.2, 3.85, 4.85, fill=LIGHT_GRAY)
    add_label(slide, "融合规则", 4.92, 1.38, 1.15)
    add_bullets(
        slide,
        [
            "SKU 行属性优先，避免 LLM 改写已确定的 gender/season/role。",
            "纯货号输入可跳过 LLM，减少成本和误判。",
            "图片高相似度命中时覆盖 gender、season、anchor_role。",
            "显式文本约束优先于图搜/LLM 覆盖。",
            "target_slots 支持按角色的正向/否定约束。",
        ],
        4.95,
        1.85,
        3.35,
        2.1,
        size=11,
    )
    add_label(slide, "输出", 4.92, 4.4, 0.9)
    add_bullets(
        slide,
        [
            "UserIntent：性别、年龄、季节、风格、场合、预算。",
            "anchor_role / target_roles：驱动后续拼套角色。",
            "anchor_attrs：结构化属性进入召回过滤。",
        ],
        4.95,
        4.85,
        3.35,
        0.85,
        size=10,
    )
    add_box(slide, 8.95, 1.2, 3.9, 4.85, fill=LIGHT_GRAY)
    add_label(slide, "评审判断", 9.12, 1.38, 1.25)
    add_bullets(
        slide,
        [
            "意图解析不是单纯 LLM，而是 SKU 属性、图向量、词典归一化和 LLM 的融合。",
            "优势：可解释、可通过配置调阈值，并能压低纯 SKU 场景的 LLM 调用。",
            "风险：图片近邻阈值、角色映射和中类归一化会级联影响召回。",
            "建议：为核心 slot 输出记录来源和置信度，用于线上问题定位和回归评估。",
        ],
        9.15,
        1.85,
        3.35,
        2.6,
        size=11,
    )
    add_text(
        slide,
        "评审关注：意图层决定召回边界，错误 slot 会同时影响 ES 过滤、"
        "Milvus expr 和拼套规则。",
        0.55,
        6.38,
        12.2,
        0.45,
        size=12,
        color=GRAY,
    )


def recall_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "1_浅色风格封面")
    add_title(slide, "搭配召回：四路候选与规则过滤")
    headers = ["召回通路", "核心函数", "候选来源", "适用场景"]
    rows = [
        [
            "image_vector",
            "recall_anchor_graph_outfits",
            "锚点 SKU -> ES outfits_by_skus_batch",
            "用户给图或明确锚点，优先找固定搭配",
        ],
        [
            "text_vector",
            "recall_text_vector_composed_outfits",
            "Milvus sku_text_vectors 按 role 召回",
            "自然语言描述，需要语义补齐搭配",
        ],
        [
            "query2es",
            "recall_query2es_composed_outfits",
            "resolve_es_query_for_role -> ES skus",
            "颜色/品类/季节等结构化约束明确",
        ],
        [
            "complementary_model",
            "recall_complementary_composed_outfits",
            "互补向量模型 + Milvus",
            "需要模型学习到的搭配兼容性",
        ],
    ]
    table = add_table(slide, headers, rows, (0.45, 1.15, 12.45, 2.65))
    for idx, width in enumerate([1.7, 3.1, 3.75, 3.9]):
        table.columns[idx].width = Inches(width)
    add_box(slide, 0.55, 4.15, 3.85, 1.8, fill=LIGHT_GRAY)
    add_label(slide, "拼套", 0.72, 4.32, 0.85)
    add_bullets(
        slide,
        [
            "compose_outfits_from_role_recall 按 target_roles 组合。",
            "order_outfit_items_by_role 保持角色展示顺序。",
            "synthetic outfit 通过 synth_* 标识来源。",
        ],
        0.75,
        4.72,
        3.35,
        0.95,
        size=10,
    )
    add_box(slide, 4.75, 4.15, 3.85, 1.8, fill=LIGHT_GRAY)
    add_label(slide, "过滤", 4.92, 4.32, 0.85)
    add_bullets(
        slide,
        [
            "category_l2_pairing 控制品类可搭配关系。",
            "color_series_pairing 控制色系兼容。",
            "outfit_conflict 过滤场景域、长度、覆盖度等冲突。",
        ],
        4.95,
        4.72,
        3.35,
        0.95,
        size=10,
    )
    add_box(slide, 8.95, 4.15, 3.9, 1.8, fill=LIGHT_GRAY)
    add_label(slide, "融合", 9.12, 4.32, 0.85)
    add_bullets(
        slide,
        [
            "merge_and_dedupe_outfits 合并多路候选。",
            "SKU 集合相同视为重复，优先保留 anchor is_master。",
            "RRF/分数用于跨通路择优。",
        ],
        9.15,
        4.72,
        3.35,
        0.95,
        size=10,
    )
    add_text(
        slide,
        "评审关注：固定搭配保证运营确定性，文本向量和 Query2ES 保证泛化；"
        "多路召回必须用规则过滤和去重控制质量。",
        0.55,
        6.38,
        12.2,
        0.45,
        size=12,
        color=GRAY,
    )


def ranking_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "1_浅色风格封面")
    add_title(slide, "粗排、LLM 排序与理由生成")
    add_box(slide, 0.55, 1.25, 3.85, 4.75, fill=LIGHT_GRAY)
    add_label(slide, "1 粗排截断", 0.72, 1.42, 1.25)
    add_bullets(
        slide,
        [
            "coarse_rank_outfits 控制进入 LLM 的候选规模。",
            "coarse_ranking_method 支持 rule/model。",
            "规则分项：source_match、intent_match、完整度、中类、色系、预算、多样性。",
            "rank_outfit_limit 控制最终返回套数。",
        ],
        0.75,
        1.9,
        3.35,
        1.7,
        size=11,
    )
    add_label(slide, "关键代码", 0.72, 4.2, 1.25)
    add_bullets(
        slide,
        [
            "ranking/outfit_ranker.py",
            "ranking/scoring.py",
            "services/outfit_recall.py",
        ],
        0.75,
        4.65,
        3.35,
        0.85,
        size=10,
    )
    add_box(slide, 4.75, 1.25, 3.85, 4.75, fill=LIGHT_GRAY)
    add_label(slide, "2 LLM 精排", 4.92, 1.42, 1.25)
    add_bullets(
        slide,
        [
            "llm_rank_outfits 调用 ranking_outfit_score.md。",
            "支持 batch / parallel 评分策略。",
            "产出 _llm_score、_llm_brief、_llm_reason。",
            "LLM 只处理粗排后的候选，避免全量候选直接进入模型。",
        ],
        4.95,
        1.9,
        3.35,
        1.7,
        size=11,
    )
    add_label(slide, "关键配置", 4.92, 4.2, 1.25)
    add_bullets(
        slide,
        [
            "models.ranking_llm: qwen3.5-flash",
            "prompt_files.ranking_outfit_score",
            "recommend.outfit_rank_weights",
        ],
        4.95,
        4.65,
        3.35,
        0.85,
        size=10,
    )
    add_box(slide, 8.95, 1.25, 3.9, 4.75, fill=LIGHT_GRAY)
    add_label(slide, "3 推荐理由", 9.12, 1.42, 1.25)
    add_bullets(
        slide,
        [
            "generate_outfit_reason_payload 统一生成理由。",
            "outfit_reason_mode 支持 outfit_only / per_item / both。",
            "reason_generation_mode=template 时可走 DphsReasonStore 话术库。",
            "最终 reason 写入 outfit_card，随搭配结果返回。",
        ],
        9.15,
        1.9,
        3.35,
        1.7,
        size=11,
    )
    add_label(slide, "关键代码", 9.12, 4.2, 1.25)
    add_bullets(
        slide,
        [
            "backend/llm_client.py",
            "prompt/reason_*.md",
            "services/card_builder.py",
        ],
        9.15,
        4.65,
        3.35,
        0.85,
        size=10,
    )
    add_text(
        slide,
        "评审关注：排序采用“规则/模型粗排 + LLM 精排 + 理由生成”的两段式，"
        "用候选截断平衡质量、成本和延迟。",
        0.55,
        6.38,
        12.2,
        0.45,
        size=12,
        color=GRAY,
    )


def responsibility_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "1_浅色风格封面")
    add_title(slide, "模块边界与关键依赖")
    headers = ["环节", "主文件 / 类", "输入", "输出"]
    rows = [
        [
            "ETL",
            "run_processed_etl.py / daily_incremental_update.py",
            "Hive CSV、运营搭配表",
            "skus.jsonl、spu_to_skus、fila_outfits.json",
        ],
        [
            "ES 索引",
            "build_fila_es_index.py / EsClient",
            "skus.jsonl、fila_outfits.json",
            "skus/outfits 索引",
        ],
        [
            "向量索引",
            "build_fila_milvus_multimodal_index.py / build_text_milvus_index.py",
            "index_images、title+色系+季节",
            "sku_vectors、sku_text_vectors",
        ],
        [
            "意图",
            "intent_engine.extract_intent / parse_user_intent",
            "文本、图片、SKU 行",
            "UserIntent、anchor_attrs、target_roles",
        ],
        [
            "召回",
            "multi_path_recall / SkuRetriever / DataFacade",
            "UserIntent、锚点 SKU、ES、Milvus",
            "多路 outfit candidates",
        ],
        [
            "排序理由",
            "llm_rank_outfits / generate_outfit_reason_payload",
            "候选搭配、用户意图、prompt",
            "排序分、理由、outfit_card",
        ],
    ]
    table = add_table(slide, headers, rows, (0.45, 1.15, 12.45, 4.15))
    for idx, width in enumerate([1.35, 4.0, 3.0, 4.1]):
        table.columns[idx].width = Inches(width)
    add_box(slide, 0.55, 5.65, 12.3, 0.85, fill=LIGHT_RED)
    add_text(
        slide,
        "边界判断：ETL 与索引构建通过 data/ 目录和索引状态解耦；线上服务"
        "通过 DataFacade 屏蔽 ES 与本地 fallback；召回和排序通过 outfit dict"
        "传递，仍需依赖字段契约保持稳定。",
        0.75,
        5.83,
        11.85,
        0.38,
        size=12,
        bold=True,
        color=DARK_RED,
        align=PP_ALIGN.CENTER,
    )


def risks_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "1_浅色风格封面")
    add_title(slide, "评审风险与改进建议")
    headers = ["风险点", "影响", "建议"]
    rows = [
        [
            "ETL 字段漂移",
            "角色、中类、色系、季节缺失会导致 ES/Milvus 召回失真",
            "增加日更字段完整率门禁，失败时阻断索引发布",
        ],
        [
            "ES 与 Milvus 不一致",
            "图向量命中 SKU 但 ES 查不到详情或搭配",
            "以 index_sync_state 为基础输出双索引一致性报告",
        ],
        [
            "意图 slot 错误",
            "错误过滤条件会放大到所有召回通路",
            "保留 slot 来源、置信度和覆盖原因，沉淀回归集",
        ],
        [
            "多路召回质量波动",
            "拼套重复、冲突或空召回影响推荐稳定性",
            "按通路统计召回量、去重率、冲突过滤率和最终入选率",
        ],
        [
            "LLM 延迟与成本",
            "精排和理由生成影响端到端响应",
            "固定候选上限、缓存理由、失败时降级到规则排序和模板理由",
        ],
    ]
    table = add_table(slide, headers, rows, (0.45, 1.2, 12.45, 4.35))
    for idx, width in enumerate([2.2, 4.45, 5.8]):
        table.columns[idx].width = Inches(width)
    add_box(slide, 0.55, 5.95, 12.3, 0.65, fill=LIGHT_RED)
    add_text(
        slide,
        "上线建议：先补齐数据质量门禁、双索引一致性校验和通路级指标；"
        "LLM 排序/理由设置明确超时与降级策略后再扩大流量。",
        0.75,
        6.12,
        11.85,
        0.3,
        size=12,
        bold=True,
        color=DARK_RED,
        align=PP_ALIGN.CENTER,
    )


def conclusion_slide(prs: Presentation) -> None:
    slide = add_slide(prs, "1_浅色风格封面")
    add_title(slide, "评审结论")
    add_box(slide, 0.55, 1.25, 5.9, 4.75, fill=LIGHT_GRAY)
    add_label(slide, "架构结论", 0.72, 1.42, 1.25)
    add_bullets(
        slide,
        [
            "数据链路完整：Hive CSV -> ETL 产物 -> ES/Milvus 在线索引。",
            "召回链路清晰：固定搭配确定性 + 语义/结构化拼套泛化。",
            "排序链路合理：规则/模型粗排控制候选，LLM 精排提升美学判断。",
            "理由链路独立：支持 LLM prompt 和模板话术两种模式。",
        ],
        0.75,
        1.9,
        5.45,
        1.75,
        size=12,
    )
    add_label(slide, "评审通过条件", 0.72, 4.2, 1.55)
    add_bullets(
        slide,
        [
            "ETL 和索引重建可重跑、可增量、可校验。",
            "线上召回通路可配置、可观测、可降级。",
            "LLM 排序和理由具备超时兜底。",
        ],
        0.75,
        4.65,
        5.45,
        0.85,
        size=11,
    )
    add_box(slide, 6.85, 1.25, 6.0, 4.75, fill=LIGHT_GRAY)
    add_label(slide, "后续优先级", 7.02, 1.42, 1.35)
    priorities = [
        ("P0", "数据门禁", "字段完整率、图片可用率、搭配关系完整率。"),
        ("P0", "索引一致性", "ES 文档数、Milvus 向量数、孤立 SKU 清理。"),
        ("P1", "通路指标", "各召回通路的候选量、命中率、入选率。"),
        ("P1", "LLM 降级", "精排失败回退规则分，理由失败回退模板。"),
        ("P2", "回归集", "按典型 query、SKU、图片样例做链路回归。"),
    ]
    for idx, (level, title, desc) in enumerate(priorities):
        y = 1.95 + idx * 0.62
        add_box(slide, 7.05, y, 0.55, 0.34, fill=RED, line=RED)
        add_text(
            slide,
            level,
            7.05,
            y + 0.04,
            0.55,
            0.2,
            size=9,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
        add_text(slide, title, 7.75, y, 1.4, 0.3, size=11, bold=True, color=RED)
        add_text(slide, desc, 9.1, y, 3.45, 0.35, size=10, color=DARK)
    add_text(
        slide,
        "结论：架构可进入评审通过路径，建议以“数据质量 + 通路指标 + "
        "LLM 降级”作为上线前重点补齐项。",
        0.55,
        6.42,
        12.2,
        0.45,
        size=12,
        bold=True,
        color=GREEN,
        align=PP_ALIGN.CENTER,
    )


def ending(prs: Presentation) -> None:
    slide = add_slide(prs, "浅色风格末页")
    add_text(
        slide,
        "感谢评审",
        0.8,
        2.8,
        11.7,
        1.0,
        size=48,
        bold=True,
        color=RED,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "FILA 穿搭推荐系统 · 推荐链路架构评审",
        0.8,
        3.85,
        11.7,
        0.5,
        size=18,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )


def main() -> int:
    if not TEMPLATE.exists():
        raise FileNotFoundError(TEMPLATE)
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)
    cover(prs)
    scope_slide(prs)
    etl_slide(prs)
    index_slide(prs)
    online_flow_slide(prs)
    intent_slide(prs)
    recall_slide(prs)
    ranking_slide(prs)
    responsibility_slide(prs)
    risks_slide(prs)
    conclusion_slide(prs)
    ending(prs)
    os.makedirs(OUTPUT.parent, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"已生成: {OUTPUT}")
    print(f"幻灯片数: {len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
