"""
生成 FILA 穿搭推荐系统架构评审 PPT
风格参考: docs/架构变更评审.pptx (浅色风格模板)
内容来源: docs/FILA+穿搭推荐系统架构.pdf + 项目代码
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

TEMPLATE = os.path.join(os.path.dirname(__file__), "..", "docs", "架构变更评审.pptx")
OUTPUT = os.path.join(os.path.dirname(__file__), "..", "docs", "FILA穿搭推荐系统架构评审_v2.pptx")

# 配色 (与参考模板一致)
RED = RGBColor(0xC0, 0x00, 0x00)        # 主色-标题强调
DARK = RGBColor(0x33, 0x33, 0x33)       # 正文深色
GRAY = RGBColor(0x59, 0x59, 0x59)       # 次要文字
LIGHT_BG = RGBColor(0xFF, 0xE1, 0xE1)   # 浅红分类底
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2) # 浅灰底
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
HEAD_BG = RGBColor(0xC0, 0x00, 0x00)    # 表头底
ALT_ROW = RGBColor(0xFB, 0xEA, 0xEA)    # 隔行底
BOX_BORDER = RGBColor(0xD0, 0xD0, 0xD0)

FONT = "微软雅黑"

# ---------- 工具函数 ----------
def remove_all_slides(prs):
    """删除模板所有 slide 及其 part，保留 master/layout（避免重复 zip 条目）"""
    sldIdLst = prs.slides._sldIdLst
    pkg = prs.part.package
    # 收集要删除的 slide part
    slide_parts_to_drop = []
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        slide_part = prs.part.related_part(rId)
        if slide_part is not None:
            slide_parts_to_drop.append(slide_part)
        # 删除 presentation part 上的 relationship
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)
    # 从 package 中移除孤立 slide part
    for sp in slide_parts_to_drop:
        try:
            pkg._parts.pop(sp.partname, None)
        except Exception:
            pass

def add_slide(prs, layout_name):
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return prs.slides.add_slide(layout)
    raise ValueError(f"layout {layout_name} not found")

def set_font(run, size=14, bold=False, color=DARK, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    # 中文字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)
    run.font.color.rgb = color

def add_title(slide, text, size=28):
    """添加左上角标题，与参考模板一致"""
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(12.4), Inches(0.7))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=True, color=RED)
    return tb

def add_text(slide, text, left, top, width, height, size=14, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color)
    return tb

def add_box(slide, left, top, width, height, fill=LIGHT_BG, line=BOX_BORDER,
            line_w=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    # 圆角调小
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh

def add_box_with_text(slide, text, left, top, width, height,
                      fill=LIGHT_BG, size=12, bold=False, color=DARK,
                      line=BOX_BORDER, align=PP_ALIGN.CENTER):
    add_box(slide, left, top, width, height, fill=fill, line=line)
    tb = add_text(slide, text, left, top, width, height, size=size,
                  bold=bold, color=color, align=align, anchor=MSO_ANCHOR.MIDDLE)
    return tb

def add_category_label(slide, text, left, top, width=1.0, height=0.34):
    """左侧分类标签 (深红底白字)"""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RED
    sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.15
    except Exception:
        pass
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_font(r, size=11, bold=True, color=WHITE)
    return sh

def add_table(slide, rows, cols, left, top, width, height):
    tbl_shape = slide.shapes.add_table(rows, cols,
                                       Inches(left), Inches(top),
                                       Inches(width), Inches(height))
    tbl = tbl_shape.table
    return tbl_shape, tbl

def style_table(tbl, header_fill=HEAD_BG, header_color=WHITE,
                body_size=12, header_size=12, col_widths=None):
    # 表头
    for j, cell in enumerate(tbl.rows[0].cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.06)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                set_font(run, size=header_size, bold=True, color=header_color)
    # 数据行
    for i in range(1, len(tbl.rows)):
        fill = ALT_ROW if i % 2 == 1 else WHITE
        for cell in tbl.rows[i].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    set_font(run, size=body_size, bold=False, color=DARK)
    if col_widths:
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(w)

def set_cell(tbl, row, col, text, size=12, bold=False, color=DARK,
             align=PP_ALIGN.LEFT):
    cell = tbl.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color)

# ---------- 构建幻灯片 ----------
prs = Presentation(TEMPLATE)
remove_all_slides(prs)

SW = prs.slide_width / 914400  # slide width in inches
SH = prs.slide_height / 914400

# ===== Slide 1: 封面 =====
s = add_slide(prs, "3_浅色风格封面")
add_text(s, "技术架构评审", Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.2),
         size=54, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, "FILA 穿搭推荐系统", Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.7),
         size=28, bold=False, color=DARK, align=PP_ALIGN.CENTER)
add_text(s, "图文对话推荐  ·  多路召回  ·  ES + Milvus 索引  ·  离线日更 ETL",
         Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.5),
         size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "2026.07", Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.4),
         size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ===== Slide 2: 产品内容 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "产品内容")
# 4 个功能板块
blocks = [
    ("用户入口", ["HTML 调试台", "SSE 对话推荐", "搭配 / 商品详情", "图片 Debug"], 0.45),
    ("在线推荐", ["意图解析", "锚点 SKU 检索", "多路搭配召回", "排序与理由"], 3.55),
    ("数据资产", ["Hive 商品表", "微导购固定搭配", "SKU/SPU 目录", "图字段选型"], 6.65),
    ("运营评测", ["批量评测", "人工 Review", "缺陷分析", "日志回放"], 9.75),
]
for title, items, left in blocks:
    add_box(s, left, 1.25, 2.9, 5.4, fill=LIGHT_GRAY)
    add_category_label(s, title, left + 0.1, 1.35, 2.7)
    for i, it in enumerate(items):
        add_box_with_text(s, it, left + 0.15, 1.95 + i * 1.1, 2.6, 0.9,
                          fill=WHITE, size=13, color=DARK, line=BOX_BORDER)
# 底部说明
add_text(s, "评审重点：覆盖从对话入口到数据资产的完整链路，重点关注推荐链路稳定性与数据日更质量。",
         0.45, 6.75, 12.4, 0.4, size=11, color=GRAY)

# ===== Slide 3: 技术架构 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "技术架构")
# 列分类标签
add_category_label(s, "前端", 0.45, 1.3, 1.0)
add_category_label(s, "应用", 0.45, 2.55, 1.0)
add_category_label(s, "模型", 0.45, 3.8, 1.0)
add_category_label(s, "存储", 0.45, 5.05, 1.0)
add_category_label(s, "工程", 0.45, 6.3, 1.0)

# 前端
front = ["Vue3 + TypeScript", "Outfits Viewer", "Image Debug Viewer", "HTML 调试台"]
for i, t in enumerate(front):
    add_box_with_text(s, t, 1.6, 1.3 + i*0.32, 2.0, 0.28,
                      fill=LIGHT_BG, size=11, color=DARK)

# 应用层
app = ["FastAPI 0.110+", "Uvicorn (单进程)", "Pydantic v2", "httpx / SSE", "PyYAML"]
for i, t in enumerate(app):
    add_box_with_text(s, t, 1.6, 2.55 + i*0.28, 2.0, 0.25,
                      fill=LIGHT_BG, size=11, color=DARK)

# 模型/算法
ml = ["Qwen3.6-27B (Vision LLM)", "qwen3.5-flash (Intent LLM)",
      "Qwen3-VL-Embedding (1024d)", "Trie 词典 (pyahocorasick)", "规则引擎"]
for i, t in enumerate(ml):
    add_box_with_text(s, t, 1.6, 3.8 + i*0.28, 2.5, 0.25,
                      fill=LIGHT_BG, size=11, color=DARK)

# 存储
store = ["Elasticsearch 7.9.3", "Milvus 2.6 (HNSW)", "Milvus Lite (IVF_FLAT)",
         "本地文件 (data/)", "LLM Gateway (OpenAI 兼容)"]
for i, t in enumerate(store):
    add_box_with_text(s, t, 1.6, 5.05 + i*0.28, 2.5, 0.25,
                      fill=LIGHT_BG, size=11, color=DARK)

# 工程
eng = ["Git", "pip + requirements.txt", "pytest", "uvicorn 脚本", "cron 日更"]
for i, t in enumerate(eng):
    add_box_with_text(s, t, 1.6, 6.3 + i*0.28, 2.0, 0.25,
                      fill=LIGHT_BG, size=11, color=DARK)

# 右侧：访问层/编排层/检索层 区块
add_box(s, 4.4, 1.3, 4.3, 1.05, fill=LIGHT_GRAY)
add_text(s, "访问层", 4.5, 1.35, 1.2, 0.3, size=11, bold=True, color=RED)
add_text(s, "REST / SSE API  ·  静态资源  ·  路由", 5.6, 1.35, 3.0, 0.6, size=11, color=DARK)

add_box(s, 4.4, 2.55, 4.3, 1.5, fill=LIGHT_GRAY)
add_text(s, "推荐编排层", 4.5, 2.6, 1.5, 0.3, size=11, bold=True, color=RED)
add_text(s, "RecommendService\n意图→召回→排序→理由→试穿\n全链路编排 + 阶段事件",
         5.7, 2.6, 2.9, 1.4, size=11, color=DARK)

add_box(s, 4.4, 4.2, 4.3, 1.35, fill=LIGHT_GRAY)
add_text(s, "检索层", 4.5, 4.25, 1.0, 0.3, size=11, bold=True, color=RED)
add_text(s, "ES Client  ·  Milvus Client\n图/文向量检索  ·  拼套规则",
         5.6, 4.25, 3.0, 1.25, size=11, color=DARK)

add_box(s, 4.4, 5.7, 4.3, 1.05, fill=LIGHT_GRAY)
add_text(s, "外部依赖", 4.5, 5.75, 1.3, 0.3, size=11, bold=True, color=RED)
add_text(s, "阿里云 ES  ·  托管 Milvus\nLLM Gateway  ·  OBS",
         5.7, 5.75, 2.9, 0.95, size=11, color=DARK)

# 右侧规范
add_box(s, 8.95, 1.3, 4.0, 5.45, fill=LIGHT_GRAY)
add_category_label(s, "技术规范", 9.05, 1.4, 3.8)
specs = [
    ("工程规范", "目录分层 / 配置 yaml 化"),
    ("代码风格", "PEP8 / type hints / mypy"),
    ("质量规范", "pytest 单测 + eval 批评"),
    ("接口规范", "REST JSON + SSE 事件流"),
    ("数据规范", "字段统一 / 增量同步状态"),
    ("Prompt 规范", "Markdown 模板 + 版本管理"),
    ("配置规范", "yaml + 环境变量覆盖"),
    ("日志规范", "trace_id + 阶段事件 + jsonl"),
]
for i, (t, d) in enumerate(specs):
    y = 1.85 + i * 0.6
    add_text(s, t, 9.1, y, 1.4, 0.28, size=11, bold=True, color=RED)
    add_text(s, d, 10.5, y, 2.4, 0.28, size=10, color=DARK)

add_text(s, "评审重点：FastAPI 单进程承载推荐全链路；模型层依赖 LLM Gateway，存储层 ES + Milvus 双模式（云/本地）。",
         0.45, 6.85, 12.4, 0.4, size=11, color=GRAY)

# ===== Slide 4: 应用架构 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "应用架构")

layers = [
    ("表现层", ["HTML 调试台", "Outfits Viewer", "Image Debug", "Eval Review"], 1.2),
    ("网关层", ["FastAPI 路由", "StaticFiles", "SSE Streaming"], 2.35),
    ("应用层", ["RecommendService", "ChatService", "CardBuilder", "TryonService"], 3.5),
    ("服务层", ["IntentEngine", "OutfitRecall", "SkuRetriever", "Ranking", "LLMClient", "EmbeddingClient"], 4.65),
    ("数据层", ["ES Client", "Milvus Client", "DataFacade", "LocalDataStore"], 5.8),
]
for name, mods, top in layers:
    add_category_label(s, name, 0.45, top, 1.1, 0.95)
    add_box(s, 1.65, top, 11.2, 0.95, fill=LIGHT_GRAY)
    n = len(mods)
    bw = 11.0 / n
    for i, m in enumerate(mods):
        add_box_with_text(s, m, 1.75 + i*bw, top + 0.12, bw - 0.1, 0.72,
                          fill=WHITE, size=11, color=DARK, line=BOX_BORDER)

# 右侧外部
add_box_with_text(s, "浏览器 / 客户端", 1.65, 6.95, 5.4, 0.35,
                  fill=LIGHT_BG, size=12, bold=True, color=RED)
add_box_with_text(s, "外部服务：LLM Gateway / 阿里云 ES / 托管 Milvus / OBS",
                  7.15, 6.95, 5.7, 0.35, fill=LIGHT_BG, size=11, color=DARK)

# ===== Slide 5: 推荐链路集成架构图 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "推荐链路集成架构图")

# 横向流程：5 个阶段
stages = [
    ("① 意图解析", "Trie 词典 + LLM\n性别/季节/风格/场合"),
    ("② 多路召回", "锚点图召回\n文本向量拼套\nQuery2ES 拼套"),
    ("③ 粗排", "规则打分截断\n性别/季节/标签"),
    ("④ LLM 排序", "美学打分\nbatch/parallel\n推荐理由生成"),
    ("⑤ 虚拟试穿", "batch_tryon\n生成试穿效果图"),
]
n = len(stages)
bw = 2.35
for i, (t, d) in enumerate(stages):
    left = 0.45 + i * (bw + 0.15)
    add_box(s, left, 1.3, bw, 1.8, fill=LIGHT_BG)
    add_text(s, t, left, 1.4, bw, 0.4, size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(s, d, left + 0.1, 1.85, bw - 0.2, 1.2, size=11, color=DARK, align=PP_ALIGN.CENTER)
    if i < n - 1:
        add_text(s, "→", left + bw, 1.9, 0.2, 0.5, size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)

# 下方：融合去重
add_box_with_text(s, "merge_and_dedupe  ·  RRF 融合 + 去重", 0.45, 3.25, 12.4, 0.45,
                  fill=LIGHT_GRAY, size=13, bold=True, color=RED, line=BOX_BORDER)

# 服务依赖区
add_box(s, 0.45, 3.9, 6.0, 2.6, fill=LIGHT_GRAY)
add_text(s, "服务编排 (RecommendService)", 0.55, 3.95, 5.8, 0.3, size=12, bold=True, color=RED)
svc = [
    "IntentEngine — Trie+LLM 意图解析",
    "ImageEmbedding — 锚点 SKU 图检索",
    "OutfitRecall — 固定搭配 + 拼套召回",
    "SkuRetriever — 单品查询 / 互补召回",
    "Ranking — 规则 / LLM 打分排序",
    "CardBuilder — 搭配卡片组装",
    "TryonService — 试穿图生成",
]
for i, t in enumerate(svc):
    add_text(s, t, 0.6, 4.3 + i*0.3, 5.8, 0.28, size=11, color=DARK)

# 外部依赖
add_box(s, 6.65, 3.9, 6.2, 2.6, fill=LIGHT_GRAY)
add_text(s, "外部依赖", 6.75, 3.95, 6.0, 0.3, size=12, bold=True, color=RED)
ext = [
    "Elasticsearch — fila-skus / fila-outfits 索引",
    "Milvus — filaskuvectors / filaskutext_vectors",
    "LLM Gateway — Qwen3.6 排序 + 推荐理由",
    "Embedding — Qwen3-VL-Embedding 1024d",
    "Trie 词典 — 性别/季节/风格/场合/角色",
    "搭配规则 — category_l2_pairing / color_series",
]
for i, t in enumerate(ext):
    add_text(s, t, 6.8, 4.3 + i*0.3, 6.0, 0.28, size=11, color=DARK)

# 输出
add_box_with_text(s, "输出：outfit_cards + 推荐理由 + 试穿效果图 (SSE 流式返回)",
                  0.45, 6.6, 12.4, 0.4, fill=LIGHT_BG, size=12, bold=True, color=RED, line=BOX_BORDER)

add_text(s, "评审重点：推荐编排仍在单进程内完成；未来高并发上线，可优先拆分召回与排序为独立服务。",
         0.45, 7.05, 12.4, 0.35, size=11, color=GRAY)

# ===== Slide 6: 数据与索引流向 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "数据与索引流向")

# ETL 流程横向
etl = [
    ("Hive 日更表", "product_* / guide_*"),
    ("Catalog ETL", "skus.jsonl / spu map"),
    ("图片选型", "display / index / tryon"),
    ("搭配构建", "fila_outfits.json"),
    ("质量校验", "catalog / image report"),
    ("索引写入", "ES + Milvus"),
]
bw = 2.0
for i, (t, d) in enumerate(etl):
    left = 0.45 + i * (bw + 0.07)
    add_box(s, left, 1.3, bw, 1.25, fill=LIGHT_BG)
    add_text(s, t, left, 1.4, bw, 0.35, size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(s, d, left + 0.05, 1.8, bw - 0.1, 0.7, size=10, color=DARK, align=PP_ALIGN.CENTER)
    if i < len(etl) - 1:
        add_text(s, "→", left + bw - 0.02, 1.7, 0.15, 0.5, size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)

# 索引清单表
add_text(s, "索引与存储清单", 0.45, 2.75, 6, 0.35, size=14, bold=True, color=RED)
tbl_shape, tbl = add_table(s, 5, 4, 0.45, 3.15, 12.4, 2.0)
set_cell(tbl, 0, 0, "类型")
set_cell(tbl, 0, 1, "名称")
set_cell(tbl, 0, 2, "用途")
set_cell(tbl, 0, 3, "说明")
rows = [
    ("ES 索引", "fila-skus", "SKU 单品索引", "商品属性 / 搜索文本 / 图 URL / 价格 / 中类 / 色系"),
    ("ES 索引", "fila-outfits", "搭配索引", "固定搭配组合（搭配 ID / 商品列表 / 来源标记）"),
    ("Milvus", "filaskuvectors", "SKU 图文向量", "1024d / COSINE / HNSW（云端）/ IVF_FLAT（本地）"),
    ("Milvus", "filaskutext_vectors", "SKU 文本向量", "1024d / COSINE / HNSW（云端）/ IVF_FLAT（本地）"),
]
for i, r in enumerate(rows, 1):
    for j, v in enumerate(r):
        set_cell(tbl, i, j, v, size=11)
style_table(tbl, body_size=11, col_widths=[1.3, 2.6, 2.6, 5.9])

# 本地文件
add_text(s, "本地文件存储", 0.45, 5.3, 6, 0.35, size=14, bold=True, color=RED)
local = [
    ("data/tables/", "Hive 日更商品原始表 CSV"),
    ("data/processed/", "离线 ETL 单品目录与款号映射"),
    ("data/preview/", "微导购固定搭配预览 JSON"),
    ("data/logs/", "索引同步状态 / 在线推荐日志 / 会话回放"),
]
for i, (d, desc) in enumerate(local):
    add_box_with_text(s, d, 0.45, 5.7 + i*0.32, 2.2, 0.28,
                      fill=LIGHT_BG, size=11, bold=True, color=RED)
    add_text(s, desc, 2.75, 5.7 + i*0.32, 4.0, 0.28, size=11, color=DARK)

# 数据流向说明
add_box(s, 7.0, 5.3, 5.85, 1.6, fill=LIGHT_GRAY)
add_text(s, "数据流向", 7.1, 5.35, 3, 0.3, size=12, bold=True, color=RED)
flow = "Hive 商品表日更 CSV → 离线 ETL 清洗 →\n搭配 JSON 构建 → 图片选型 → 数据校验 →\nES 索引写入 + Milvus 向量写入（支持增量）\n增量同步状态由 index_sync_state 管理"
add_text(s, flow, 7.1, 5.65, 5.65, 1.25, size=11, color=DARK)

add_text(s, "评审重点：ES 7.9.3 内网集群 + Milvus 云端/本地双模式；增量同步保证日更仅写入变更数据。",
         0.45, 7.05, 12.4, 0.35, size=11, color=GRAY)

# ===== Slide 7: 集成架构关系表 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "集成架构关系表")

tbl_shape, tbl = add_table(s, 9, 5, 0.45, 1.3, 12.4, 4.8)
headers = ["本系统模块", "集成对象", "协议", "用途", "方向"]
for j, h in enumerate(headers):
    set_cell(tbl, 0, j, h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rows = [
    ("RecommendService", "LLM Gateway", "HTTPS / OpenAI 兼容", "意图提取 / 排序打分 / 推荐理由", "出"),
    ("EmbeddingClient", "LLM Gateway", "HTTPS / OpenAI 兼容", "Qwen3-VL-Embedding 1024d 向量", "出"),
    ("ESClient", "阿里云 ES 集群", "REST (9200)", "fila-skus / fila-outfits 检索", "出"),
    ("MilvusClient", "托管 Milvus", "gRPC (19530)", "图文向量 / 文本向量检索", "出"),
    ("TryonService", "LLM Gateway", "HTTPS / OpenAI 兼容", "虚拟试穿图生成", "出"),
    ("DataFacade", "本地文件", "文件 IO", "data/processed · preview · logs", "内"),
    ("前端", "FastAPI", "HTTP / SSE", "对话推荐 / 搭配查询 / 商品详情", "入"),
    ("ETL 脚本", "Hive / ES / Milvus", "SQL / REST / gRPC", "日更下载 → 索引写入", "内"),
]
for i, r in enumerate(rows, 1):
    for j, v in enumerate(r):
        set_cell(tbl, i, j, v, size=11)
style_table(tbl, body_size=11, col_widths=[2.6, 2.6, 2.6, 3.0, 1.6])

add_text(s, "评审重点：服务边界清晰，推荐编排仍在单进程内完成；外部依赖均通过 Client 封装，便于Mock与降级。",
         0.45, 6.3, 12.4, 0.4, size=11, color=GRAY)

# ===== Slide 8: 新增资源清单 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "新增资源清单")

tbl_shape, tbl = add_table(s, 9, 5, 0.45, 1.3, 12.4, 4.5)
headers = ["资源类型", "名称 / 标识", "规格 / 配置", "用途", "所属"]
for j, h in enumerate(headers):
    set_cell(tbl, 0, j, h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rows = [
    ("ES 索引", "fila-skus", "7.9.3 / 分片 1+1", "SKU 单品检索", "内网 ES 集群"),
    ("ES 索引", "fila-outfits", "7.9.3 / 分片 1+1", "固定搭配检索", "内网 ES 集群"),
    ("Milvus Collection", "filaskuvectors", "1024d / HNSW / COSINE", "SKU 图文向量", "阿里云托管"),
    ("Milvus Collection", "filaskutext_vectors", "1024d / HNSW / COSINE", "SKU 文本向量", "阿里云托管"),
    ("LLM 模型", "qwen3.6-27b", "Vision LLM", "排序 / 推荐理由 / 试穿", "LLM Gateway"),
    ("LLM 模型", "qwen3.5-flash", "Intent LLM", "意图解析 fallback", "LLM Gateway"),
    ("Embedding", "Qwen3-VL-Embedding", "1024 维", "图文向量化", "LLM Gateway"),
    ("OBS 桶", "fila-tryon-images", "标准桶", "试穿图缓存", "华为云 OBS"),
]
for i, r in enumerate(rows, 1):
    for j, v in enumerate(r):
        set_cell(tbl, i, j, v, size=11)
style_table(tbl, body_size=11, col_widths=[2.2, 2.8, 2.6, 2.8, 2.0])

add_text(s, "评审重点：复用现有 ES / Milvus / LLM Gateway 资源，新增仅索引与 Collection 级别对象。",
         0.45, 6.0, 12.4, 0.4, size=11, color=GRAY)

# ===== Slide 9: K8S 资源清单 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "K8S 资源清单")

add_text(s, "当前部署形态：FastAPI 单进程服务（uvicorn），暂未容器化；以下为规划目标。",
         0.45, 1.2, 12.4, 0.4, size=12, color=GRAY)

tbl_shape, tbl = add_table(s, 8, 5, 0.45, 1.75, 12.4, 4.2)
headers = ["K8s 对象", "名称", "规格", "数量", "说明"]
for j, h in enumerate(headers):
    set_cell(tbl, 0, j, h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rows = [
    ("Deployment", "fila-rec", "2 vCPU / 4Gi", "2", "FastAPI + Uvicorn 单进程"),
    ("Service", "fila-rec-svc", "ClusterIP", "1", "内部调用入口"),
    ("Ingress", "fila-rec-ingress", "HTTPS / 域名", "1", "对外 REST + SSE"),
    ("ConfigMap", "fila-rec-config", "config.yaml", "1", "运行时配置"),
    ("Secret", "fila-rec-secret", "API Key / 凭证", "1", "LLM / ES / Milvus 认证"),
    ("HPA", "fila-rec-hpa", "CPU 70%", "1", "自动扩缩容 2~6"),
    ("CronJob", "fila-daily-etl", "2 vCPU / 2Gi", "1", "日更 ETL 编排"),
]
for i, r in enumerate(rows, 1):
    for j, v in enumerate(r):
        set_cell(tbl, i, j, v, size=11)
style_table(tbl, body_size=11, col_widths=[2.2, 2.6, 2.6, 1.4, 3.6])

add_text(s, "评审重点：当前为单进程裸机部署，建议容器化并配置 HPA；日更 ETL 通过 CronJob 编排。",
         0.45, 6.1, 12.4, 0.4, size=11, color=GRAY)

# ===== Slide 10: 监控与可观测性 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "监控与可观测性方案")

# 已具备
add_box(s, 0.45, 1.3, 6.0, 4.8, fill=LIGHT_GRAY)
add_text(s, "✓ 已具备", 0.55, 1.4, 5.8, 0.4, size=16, bold=True, color=RGBColor(0x2E,0x8B,0x57))
have = [
    ("trace_id", "全链路追踪标识贯穿 SSE 流"),
    ("SSE 阶段事件", "意图/召回/排序/试穿进度实时回传"),
    ("召回 IO 日志", "三路召回候选与分数落盘"),
    ("jsonl_logger", "结构化会话日志按行写入"),
    ("缺陷分析日志", "eval/review 缺陷归类与回放"),
    ("日志回放", "会话级 replay 支持复现"),
    ("search_debug", "ANN/ES 检索调试台在线"),
]
for i, (t, d) in enumerate(have):
    add_text(s, t, 0.6, 1.9 + i*0.6, 1.8, 0.3, size=12, bold=True, color=RED)
    add_text(s, d, 2.5, 1.9 + i*0.6, 3.9, 0.55, size=11, color=DARK)

# 待补齐
add_box(s, 6.65, 1.3, 6.2, 4.8, fill=LIGHT_GRAY)
add_text(s, "✗ 待补齐", 6.75, 1.4, 6.0, 0.4, size=16, bold=True, color=RGBColor(0xB0,0x30,0x30))
todo = [
    ("统一指标上报", "Prometheus / 自定义 metrics"),
    ("告警阈值", "LLM 超时 / 召回空 / ES 异常"),
    ("大盘看板", "Grafana 推荐链路 QPS/延时"),
    ("SLA 监控", "端到端 P95 / 错误率"),
    ("资源水位", "ES / Milvus / LLM 配额"),
    ("数据质量监控", "日更 ETL 成功率 / 缺字段"),
    ("Tryon 监控", "试穿图生成成功率 / 延时"),
]
for i, (t, d) in enumerate(todo):
    add_text(s, t, 6.8, 1.9 + i*0.6, 1.9, 0.3, size=12, bold=True, color=RED)
    add_text(s, d, 8.8, 1.9 + i*0.6, 4.0, 0.55, size=11, color=DARK)

add_text(s, "评审重点：可观测性基础已具备 trace_id 与阶段事件；需补齐指标上报、告警阈值与 Grafana 大盘。",
         0.45, 6.3, 12.4, 0.4, size=11, color=GRAY)

# ===== Slide 11: 安全方案 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "安全方案")

secs = [
    ("认证鉴权", "API Key 校验（.env 注入）", "LLM Gateway / ES / Milvus 凭证不入仓"),
    ("传输加密", "HTTPS / TLS", "对外接口走 Ingress TLS；内网 ES/Milvus 走 VPC"),
    ("输入校验", "Pydantic 模型校验", "所有 REST 入参强类型校验，拒绝非法 payload"),
    ("注入防护", "ES Query 参数化 / Trie 词典", "用户输入不直接拼 ES DSL；LLM prompt 模板化"),
    ("敏感数据", "图片 URL / 商品数据脱敏", "日志不输出完整图片 URL 与价格明细"),
    ("限流降级", "LLM 超时 + 召回 fallback", "LLM 3s 超时回退 Trie；召回空返回兜底搭配"),
    ("配置安全", "yaml + 环境变量覆盖", "密钥通过 env / Secret 注入，config.yaml 入仓"),
    ("审计追溯", "trace_id + jsonl 会话日志", "全链路可追溯，支持会话回放与缺陷定位"),
]
tbl_shape, tbl = add_table(s, len(secs)+1, 3, 0.45, 1.3, 12.4, 5.2)
headers = ["维度", "措施", "说明"]
for j, h in enumerate(headers):
    set_cell(tbl, 0, j, h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for i, (a, b, c) in enumerate(secs, 1):
    set_cell(tbl, i, 0, a, size=11, bold=True, color=RED)
    set_cell(tbl, i, 1, b, size=11)
    set_cell(tbl, i, 2, c, size=11)
style_table(tbl, body_size=11, col_widths=[2.0, 4.2, 6.2])

add_text(s, "评审重点：密钥全部走环境变量；ES/LLM 输入参数化；LLM 超时回退保证可用性。",
         0.45, 6.7, 12.4, 0.4, size=11, color=GRAY)

# ===== Slide 12: 备份恢复与评审结论 =====
s = add_slide(prs, "1_浅色风格封面")
add_title(s, "备份恢复与评审结论")

# 左：备份恢复
add_box(s, 0.45, 1.3, 6.0, 4.0, fill=LIGHT_GRAY)
add_text(s, "备份恢复", 0.55, 1.4, 5.8, 0.4, size=16, bold=True, color=RED)
backup = [
    ("ES 索引", "支持全量重建 + 增量更新"),
    ("Milvus 向量", "云端快照 + 本地 Milvus Lite"),
    ("本地数据", "data/ 目录版本管理 (Git)"),
    ("搭配 JSON", "微导购源表 + 构建脚本可重跑"),
    ("会话日志", "jsonl 按日归档，支持回放"),
    ("配置文件", "config.yaml 入 Git；env 走 Secret"),
    ("恢复演练", "建议每季度全链路重建演练"),
]
for i, (t, d) in enumerate(backup):
    add_text(s, t, 0.6, 1.9 + i*0.45, 1.6, 0.3, size=12, bold=True, color=RED)
    add_text(s, d, 2.3, 1.9 + i*0.45, 4.1, 0.4, size=11, color=DARK)

# 右：评审结论
add_box(s, 6.65, 1.3, 6.2, 4.0, fill=LIGHT_GRAY)
add_text(s, "评审结论", 6.75, 1.4, 6.0, 0.4, size=16, bold=True, color=RED)
concl = [
    ("架构合理性", "FastAPI 单进程承载全链路，符合当前流量；模块边界清晰"),
    ("技术选型", "ES + Milvus + LLM Gateway 组合成熟，双模式支持本地调试"),
    ("数据质量", "日更 ETL + 增量同步状态机保证索引一致性"),
    ("可观测性", "trace_id + 阶段事件基础好，需补齐指标与告警"),
    ("安全合规", "密钥 env 化 + 输入校验 + LLM 超时降级"),
    ("演进建议", "高并发时优先拆分召回/排序为独立服务 + 容器化 HPA"),
    ("上线建议", "补齐监控大盘与告警阈值后可上线"),
]
for i, (t, d) in enumerate(concl):
    add_text(s, t, 6.8, 1.9 + i*0.45, 1.6, 0.3, size=12, bold=True, color=RED)
    add_text(s, d, 8.5, 1.9 + i*0.45, 4.3, 0.55, size=11, color=DARK)

# 待跟进事项
add_text(s, "待跟进事项", 0.45, 5.5, 6, 0.35, size=14, bold=True, color=RED)
items = [
    "1. Prometheus 指标接入 + Grafana 推荐链路大盘",
    "2. 告警阈值定义（LLM 超时 / 召回空 / ES 异常 / 日更失败）",
    "3. 容器化 Dockerfile + K8s 部署清单落地",
    "4. 季度全链路重建演练纳入运维流程",
]
for i, t in enumerate(items):
    add_text(s, t, 0.45, 5.9 + i*0.32, 12.4, 0.3, size=11, color=DARK)

# ===== Slide 13: 末页 =====
s = add_slide(prs, "浅色风格末页")
add_text(s, "感谢评审", Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.0),
         size=48, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, "FILA 穿搭推荐系统  ·  技术架构评审",
         Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.5),
         size=18, color=GRAY, align=PP_ALIGN.CENTER)

# 保存
prs.save(OUTPUT)
print(f"✓ 已生成: {OUTPUT}")
print(f"  幻灯片数: {len(prs.slides)}")
